from __future__ import annotations

from pathlib import Path
from shutil import copy2

from PIL import Image, ImageDraw, ImageFilter, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path("talk/2026-04-02-live/2026-04-02-live-deck.pptx")
BACKUP = Path("talk/2026-04-02-live/2026-04-02-live-deck.original.pptx")
LOGO = Path("assets/logo.png")
LOGO_WIDE = Path("assets/logo-text-colorful.png")
PLUGINS_IMG = Path("assets/weapp-tw-plugins.png")
FRAMEWORKS_IMG = Path("assets/weapp-tw-frameworks.png")
RELEASE_IMG = Path("website/blog/2025/3/assets/v4-release.png")
LAYER_IMG = Path("website/docs/quick-start/v4/tailwindcss-v4-uniapp-layer.png")
CREATE_PROJECT_IMG = Path("website/static/img/create-project.png")
DEMO_LOGO = Path("demo/uni-app-tailwindcss-v4/src/static/logo.png")
GENERATED_DIR = Path("talk/2026-04-02-live/assets/generated")
SHOWCASE_COLLAGE = GENERATED_DIR / "showcase-collage.png"
BG_GRADIENT = GENERATED_DIR / "bg-gradient.png"
PANEL_GRADIENT = GENERATED_DIR / "panel-gradient.png"
PANEL_GRADIENT_ALT = GENERATED_DIR / "panel-gradient-alt.png"
ACCENT_GRADIENT = GENERATED_DIR / "accent-gradient.png"

INK = RGBColor(10, 15, 27)
SURFACE = RGBColor(19, 28, 46)
SURFACE_2 = RGBColor(30, 41, 59)
TEXT = RGBColor(244, 248, 255)
MUTED = RGBColor(148, 163, 184)
CYAN = RGBColor(34, 211, 238)
CYAN_2 = RGBColor(14, 165, 233)
EMERALD = RGBColor(52, 211, 153)
AMBER = RGBColor(251, 191, 36)
ROSE = RGBColor(251, 113, 133)
WHITE = RGBColor(255, 255, 255)
PANEL_BORDER = RGBColor(36, 56, 92)
PANEL_INNER = RGBColor(15, 23, 42)
PANEL_TAG = RGBColor(30, 41, 59)

POSTINSTALL_SNIPPET = """{
  "scripts": {
    "dev": "npm run dev:mp-weixin",
    "build:mp-weixin": "node ../../scripts/uni-build-guard.mjs build -p mp-weixin",
    "postinstall": "weapp-tw patch"
  }
}"""

VITE_SNIPPET = """plugins: [
  uni(),
  tailwindcss(),
  UnifiedViteWeappTailwindcssPlugin({
    rem2rpx: true,
    cssEntries: [
      path.resolve(__dirname, "src/main.css"),
      path.resolve(__dirname, "src/common.css"),
    ],
  }),
]"""

MAIN_CSS_SNIPPET = """@import "weapp-tailwindcss";
@config "../tailwind.config.js";

@theme {
  --color-neutral-1B: #1b1b1b;
  --color-midnight: #121063;
  --color-tahiti: #3ab7bf;
  --color-bermuda: #78dcca;
}"""

INDEX_SNIPPET = """<view :class="className" class="aspect-[calc(4*3+1)/3]">
  {{ className }}
</view>

const className = ref(
  'bg-[#0000ff] text-[45rpx] text-white'
)"""

PROMPT_SNIPPET = """请帮我生成一组适合小程序页面的 Tailwind 类名，不要写 CSS。
目标是一个渐变卡片：
- 大圆角
- 蓝青色渐变背景
- 柔和阴影
- 标题 32rpx，加粗
- 副标题 24rpx，透明度略低
- 底部一个白底深色字按钮"""

SKILL_PROMPT_SNIPPET = """我现在是 uni-app cli vue3 vite 项目，目标端是微信小程序 + H5。
请按 weapp-tailwindcss skill 给我最小可用配置。
输出需要包含：
1. 安装命令
2. 完整配置文件
3. 页面示例
4. 验证步骤
5. 回滚方案"""

SHOWCASE_IMAGES = [
    Path("website/static/img/showcase/001.区白白/02-1.png"),
    Path("website/static/img/showcase/005.安心校友圈/02-安心-首页.jpg"),
    Path("website/static/img/showcase/010.小而美工具/02-image.png"),
    Path("website/static/img/showcase/011.元气番茄钟/01-image.jpg"),
]


SLIDES = [
    {
        "layout": "title",
        "kicker": "2026-04-02 LIVE",
        "title": "小程序还能这么写？",
        "subtitle": "AI 出样式，Tailwind 跑全端",
        "body": "一句话描述界面 -> AI 生成 Tailwind 类名 -> weapp-tailwindcss 自动转译 -> 小程序直接预览",
        "chips": ["痛点重述", "AI 出 UI", "Tailwind 语言层", "工程化落地"],
    },
    {
        "layout": "questions",
        "title": "灵魂拷问",
        "subtitle": "你到底花了多少时间在写样式？",
        "items": [
            "上周写了多少行样式？",
            "花了多少时间调 margin / padding / radius / shadow？",
            "你是在写业务，还是在做像素搬运？",
        ],
        "note": "小程序场景还要额外承担平台差异、单位适配和多端一致性。",
    },
    {
        "layout": "grid4",
        "title": "为什么小程序样式开发一直慢",
        "subtitle": "不是你不努力，是旧范式本来就低效",
        "items": ["表达成本高", "试错成本高", "跨端一致性差", "反馈链路长"],
    },
    {
        "layout": "compare",
        "title": "旧范式 vs 新范式",
        "subtitle": "从“手写属性”切到“表达意图”",
        "left_title": "旧范式",
        "left_items": ["先想 CSS 属性", "一条一条试", "人肉反复试错"],
        "right_title": "新范式",
        "right_items": ["先说清界面意图", "AI 给第一版结果", "Tailwind 做受约束表达"],
    },
    {
        "layout": "prompt",
        "title": "先看结果",
        "subtitle": "一句话，AI 出一版卡片样式",
        "prompt": "做一个渐变卡片，大圆角，柔和阴影，标题 32rpx，副标题 24rpx，底部一个高亮 CTA。不要写传统 CSS，直接给我 Tailwind 类名。",
        "bullets": ["渐变卡片", "大圆角", "柔和阴影", "标题 / 副标题 / CTA"],
    },
    {
        "layout": "grid4",
        "title": "为什么是 Tailwind，不是直接写 CSS",
        "subtitle": "Tailwind 是更适合 AI 的样式语言",
        "items": ["可枚举", "可组合", "可约束", "更稳定"],
    },
    {
        "layout": "flow3",
        "title": "三层模型",
        "subtitle": "意图 -> 原子类 -> 小程序适配",
        "items": ["你：描述界面意图", "AI：生成 Tailwind 类名", "weapp-tailwindcss：翻译成小程序结果"],
        "note": "AI 不直接对接小程序，它先对接 Tailwind 这层表达语言。",
    },
    {
        "layout": "stack",
        "title": "weapp-tailwindcss 在做什么",
        "subtitle": "不只是插件，而是一条工程链路",
        "items": [
            "Tailwind 小程序全方面解决方案",
            "多构建工具基底支持",
            "多框架、多端与运行时适配",
            "把“能写 Tailwind”升级成“能稳定运行”",
        ],
    },
    {
        "layout": "chips",
        "title": "支持范围",
        "subtitle": "覆盖主流小程序开发路线",
        "items": ["webpack", "vite", "rspack", "rollup", "rolldown", "gulp"],
        "note": "支持的是这些构建工具基底，以及其上游框架生态。",
    },
    {
        "layout": "repo",
        "title": "仓库不是 PPT 工程",
        "subtitle": "它同时有 demo、templates、Skill、E2E、benchmark",
        "items": ["demo/*", "templates.jsonc", "skills/weapp-tailwindcss", "watch / HMR 回归", "framework benchmark"],
    },
    {
        "layout": "focus",
        "title": "主 Demo 选型",
        "subtitle": "demo/uni-app-tailwindcss-v4",
        "items": ["直观", "稳定", "贴近日常业务", "适合直播"],
        "note": "现场适合展示 patch、vite 插件、样式入口和页面类名。",
    },
    {
        "layout": "config",
        "title": "关键配置 1",
        "subtitle": 'postinstall: "weapp-tw patch"',
        "items": ["patch 链路不能漏", "任意值问题常和它有关", "JS / TS 字符串 class 也常和它有关"],
    },
    {
        "layout": "config",
        "title": "关键配置 2",
        "subtitle": "vite.config.ts",
        "items": ["uni()", "Tailwind 插件", "UnifiedViteWeappTailwindcssPlugin"],
    },
    {
        "layout": "config",
        "title": "关键配置 3",
        "subtitle": "src/main.css",
        "items": ["Tailwind 入口", "主题变量", "配置承载点"],
    },
    {
        "layout": "flow4",
        "title": "AI Skill",
        "subtitle": "让 AI 不只是写代码，而是按流程做事",
        "items": ["先收集上下文", "再输出配置", "再给验证步骤", "再给回滚方案"],
    },
    {
        "layout": "chips",
        "title": "Skill 先问什么",
        "subtitle": "最小上下文决定方案质量",
        "items": ["框架", "构建器", "目标端", "Tailwind 版本", "包管理器"],
        "note": "上下文不完整，AI 就会四处漏风。",
    },
    {
        "layout": "deliverables",
        "title": "Skill 输出什么",
        "subtitle": "不是一段代码，而是一套落地结果",
        "items": ["修改文件清单", "可复制配置", "安装命令", "验证步骤", "回滚方案"],
    },
    {
        "layout": "repo",
        "title": "进阶技巧",
        "subtitle": "这些地方最容易踩坑",
        "items": ["任意值与 rpx", "动态 class 要枚举", "space-y / space-x", "twMerge / cva / cn", "uni-app x"],
    },
    {
        "layout": "metrics",
        "title": "工程信号",
        "subtitle": "从“能跑”走向“可验证”",
        "items": ["HMR 报告", "framework benchmark", "用数据讨论体验"],
        "note": "重点不是某框架必赢，而是这套生态已经把问题工程化。",
    },
    {
        "layout": "closing",
        "title": "别再把样式开发当体力活",
        "subtitle": "从下一个页面开始试",
        "items": ["AI 写意图", "Tailwind 做语言", "weapp-tailwindcss 负责翻译", "先跑一个 demo，再用一句真实需求试一次"],
    },
]


def add_textbox(slide, left, top, width, height, text, *, size=24, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.name = "PingFang SC"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_card(slide, left, top, width, height, fill, radius=MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(radius, left, top, width, height)
    if isinstance(fill, Path):
        shape.fill.background()
        shape.line.color.rgb = PANEL_BORDER
        inner_w = width - Inches(0.04)
        inner_h = height - Inches(0.04)
        slide.shapes.add_picture(str(fill), left + Inches(0.02), top + Inches(0.02), width=inner_w, height=inner_h)
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.color.rgb = fill
    return shape


def add_picture(slide, path: Path, left, top, width=None, height=None):
    if path.exists():
        if width is not None and height is not None:
            return slide.shapes.add_picture(str(path), left, top, width=width, height=height)
        if width is not None:
            return slide.shapes.add_picture(str(path), left, top, width=width)
        if height is not None:
            return slide.shapes.add_picture(str(path), left, top, height=height)
        return slide.shapes.add_picture(str(path), left, top)
    return None


def add_code_block(slide, left, top, width, height, title, code):
    add_card(slide, left, top, width, height, PANEL_GRADIENT)
    add_textbox(slide, left + Inches(0.18), top + Inches(0.12), width - Inches(0.4), Inches(0.22), title, size=11, color=CYAN, bold=True)
    box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.42), width - Inches(0.4), height - Inches(0.55))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code
    run = p.runs[0]
    run.font.name = "Menlo"
    run.font.size = Pt(9.5)
    run.font.color.rgb = RGBColor(226, 232, 240)
    return box


def add_panel_frame(slide, left, top, width, height, title=None):
    outer = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    outer.fill.background()
    outer.line.color.rgb = PANEL_BORDER
    slide.shapes.add_picture(str(PANEL_GRADIENT), left + Inches(0.02), top + Inches(0.02), width=width - Inches(0.04), height=height - Inches(0.04))
    if title:
        add_textbox(slide, left + Inches(0.18), top + Inches(0.1), width - Inches(0.36), Inches(0.18), title, size=11, color=CYAN, bold=True)
    return outer


def add_footer_note(slide, left, top, width, text):
    slide.shapes.add_picture(str(ACCENT_GRADIENT), left, top, width=width, height=Inches(0.34))
    add_textbox(slide, left + Inches(0.18), top + Inches(0.09), width - Inches(0.36), Inches(0.14), text, size=11.2, color=TEXT)


def add_tag(slide, left, top, width, text):
    slide.shapes.add_picture(str(PANEL_GRADIENT_ALT), left, top, width=width, height=Inches(0.38))
    add_textbox(slide, left + Inches(0.12), top + Inches(0.1), width - Inches(0.24), Inches(0.14), text, size=11.2, color=CYAN, bold=True, align=PP_ALIGN.CENTER)


def ensure_generated_assets():
    GENERATED_DIR.mkdir(parents=True, exist_ok=True)
    bg = Image.new("RGB", (1600, 900), (8, 12, 24))
    draw = ImageDraw.Draw(bg)
    for y in range(900):
        ratio = y / 899
        r = int(8 + (16 - 8) * ratio)
        g = int(12 + (30 - 12) * ratio)
        b = int(24 + (56 - 24) * ratio)
        draw.line((0, y, 1600, y), fill=(r, g, b))
    glow = Image.new("RGBA", (1600, 900), (0, 0, 0, 0))
    gd = ImageDraw.Draw(glow)
    gd.ellipse((980, 20, 1540, 580), fill=(0, 180, 255, 90))
    gd.ellipse((80, 540, 580, 980), fill=(0, 255, 180, 60))
    glow = glow.filter(ImageFilter.GaussianBlur(70))
    bg = Image.alpha_composite(bg.convert("RGBA"), glow)
    bg.convert("RGB").save(BG_GRADIENT)

    img = Image.new("RGB", (1200, 700), (15, 23, 42))
    draw = ImageDraw.Draw(img)
    for y in range(700):
        ratio = y / 699
        color = (
            int(18 + 10 * ratio),
            int(28 + 25 * ratio),
            int(46 + 55 * ratio),
        )
        draw.line((0, y, 1200, y), fill=color)
    glow = Image.new("RGBA", (1200, 700), (0, 0, 0, 0))
    gd = ImageDraw.Draw(glow)
    gd.ellipse((650, -80, 1180, 380), fill=(0, 200, 255, 80))
    gd.ellipse((-120, 350, 420, 820), fill=(0, 255, 180, 55))
    glow = glow.filter(ImageFilter.GaussianBlur(55))
    img = Image.alpha_composite(img.convert("RGBA"), glow)
    img.convert("RGB").save(PANEL_GRADIENT)

    img = Image.new("RGB", (900, 220), (20, 30, 52))
    draw = ImageDraw.Draw(img)
    for x in range(900):
        ratio = x / 899
        color = (
            int(22 + (0 - 22) * ratio),
            int(38 + (170 - 38) * ratio),
            int(68 + (233 - 68) * ratio),
        )
        draw.line((x, 0, x, 220), fill=color)
    glow = Image.new("RGBA", (900, 220), (0, 0, 0, 0))
    gd = ImageDraw.Draw(glow)
    gd.ellipse((420, -60, 900, 240), fill=(0, 220, 255, 95))
    gd.ellipse((-180, 60, 260, 300), fill=(0, 255, 180, 65))
    glow = glow.filter(ImageFilter.GaussianBlur(30))
    img = Image.alpha_composite(img.convert("RGBA"), glow)
    img.convert("RGB").save(PANEL_GRADIENT_ALT)

    img = Image.new("RGB", (1200, 120), (8, 80, 94))
    draw = ImageDraw.Draw(img)
    for x in range(1200):
        ratio = x / 1199
        color = (
            int(6 + (14 - 6) * ratio),
            int(119 + (165 - 119) * ratio),
            int(98 + (233 - 98) * ratio),
        )
        draw.line((x, 0, x, 120), fill=color)
    glow = Image.new("RGBA", (1200, 120), (0, 0, 0, 0))
    gd = ImageDraw.Draw(glow)
    gd.ellipse((720, -30, 1240, 150), fill=(0, 220, 255, 100))
    glow = glow.filter(ImageFilter.GaussianBlur(26))
    img = Image.alpha_composite(img.convert("RGBA"), glow)
    img.convert("RGB").save(ACCENT_GRADIENT)
    if SHOWCASE_COLLAGE.exists():
        return
    canvas = Image.new("RGB", (1600, 900), (10, 15, 27))
    slots = [
        (40, 40, 340, 380),
        (420, 40, 340, 380),
        (800, 40, 760, 380),
        (40, 460, 1520, 400),
    ]
    for src, slot in zip(SHOWCASE_IMAGES, slots):
        if not src.exists():
            continue
        img = Image.open(src).convert("RGB")
        fitted = ImageOps.fit(img, (slot[2], slot[3]), method=Image.Resampling.LANCZOS)
        canvas.paste(fitted, (slot[0], slot[1]))
    canvas.save(SHOWCASE_COLLAGE)


def apply_background(slide, page_num):
    slide.shapes.add_picture(str(BG_GRADIENT), 0, 0, width=Inches(13.333), height=Inches(7.5))
    slide.shapes.add_picture(str(ACCENT_GRADIENT), 0, 0, width=Inches(0.12), height=Inches(7.5))
    add_textbox(slide, Inches(12.3), Inches(7.02), Inches(0.6), Inches(0.24), f"{page_num:02d}", size=11, color=MUTED, align=PP_ALIGN.RIGHT)


def add_header(slide, title, subtitle):
    add_textbox(slide, Inches(0.7), Inches(0.48), Inches(8.8), Inches(0.6), title, size=27, bold=True)
    add_textbox(slide, Inches(0.72), Inches(1.08), Inches(9.6), Inches(0.36), subtitle, size=13, color=MUTED)


def build(prs: Presentation):
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for idx, spec in enumerate(SLIDES, start=1):
        slide = prs.slides.add_slide(blank)
        apply_background(slide, idx)
        layout = spec["layout"]
        if layout == "title":
            add_panel_frame(slide, Inches(7.92), Inches(0.72), Inches(4.6), Inches(0.96), "Brand")
            add_picture(slide, LOGO_WIDE, Inches(8.22), Inches(0.96), width=Inches(3.95), height=Inches(0.64))
            add_textbox(slide, Inches(0.9), Inches(0.72), Inches(2.2), Inches(0.3), spec["kicker"], size=11, color=CYAN, bold=True)
            add_textbox(slide, Inches(0.88), Inches(1.38), Inches(7.4), Inches(1.4), spec["title"], size=29, bold=True)
            add_textbox(slide, Inches(0.88), Inches(2.34), Inches(6.8), Inches(0.6), spec["subtitle"], size=22, color=WHITE)
            add_card(slide, Inches(0.88), Inches(3.22), Inches(5.95), Inches(1.14), PANEL_GRADIENT)
            add_textbox(slide, Inches(1.12), Inches(3.5), Inches(5.5), Inches(0.6), spec["body"], size=14)
            for i, chip in enumerate(spec["chips"]):
                left = Inches(0.9 + (i % 2) * 2.25)
                top = Inches(5.18 + (i // 2) * 0.74)
                add_card(slide, left, top, Inches(2.0), Inches(0.5), PANEL_GRADIENT_ALT)
                add_textbox(slide, left + Inches(0.16), top + Inches(0.12), Inches(1.7), Inches(0.2), chip, size=11)
            hero = add_panel_frame(slide, Inches(8.12), Inches(1.84), Inches(4.45), Inches(4.3), "Pipeline")
            add_textbox(slide, Inches(8.55), Inches(1.7), Inches(3.5), Inches(0.4), "意图 -> 原子类 -> 运行结果", size=15, color=CYAN, bold=True)
            for i, word in enumerate(["Prompt", "Tailwind", "weapp-tw", "Preview"]):
                y = 2.4 + i * 0.8
                add_card(slide, Inches(8.58), Inches(y), Inches(3.3), Inches(0.46), ACCENT_GRADIENT if i % 2 == 0 else PANEL_GRADIENT_ALT)
                add_textbox(slide, Inches(8.76), Inches(y + 0.11), Inches(2.9), Inches(0.2), word, size=12, color=TEXT, bold=True)
            continue

        if layout != "closing":
            add_header(slide, spec["title"], spec["subtitle"])

        if layout == "questions":
            add_textbox(slide, Inches(0.95), Inches(1.62), Inches(6.7), Inches(1.1), "你不是在写样式。\n你是在交时间税。", size=29, bold=True)
            for i, item in enumerate(spec["items"]):
                top = Inches(3.36 + i * 0.82)
                add_card(slide, Inches(0.95), top, Inches(6.05), Inches(0.52), PANEL_GRADIENT)
                add_textbox(slide, Inches(1.14), top + Inches(0.12), Inches(5.6), Inches(0.18), item, size=12.8, bold=True)
            add_card(slide, Inches(8.0), Inches(1.92), Inches(4.0), Inches(3.9), PANEL_GRADIENT)
            add_textbox(slide, Inches(8.28), Inches(2.02), Inches(3.1), Inches(0.72), "时间税", size=28, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
            add_textbox(slide, Inches(8.34), Inches(2.25), Inches(3.0), Inches(0.32), "额外负担", size=14, color=CYAN, bold=True)
            for i, word in enumerate(["平台差异", "单位适配", "多端一致性"]):
                top = Inches(2.95 + i * 0.82)
                add_card(slide, Inches(8.36), top, Inches(2.95), Inches(0.48), PANEL_GRADIENT_ALT if i != 1 else ACCENT_GRADIENT)
                add_textbox(slide, Inches(8.5), top + Inches(0.12), Inches(2.65), Inches(0.16), word, size=12.5, bold=True, align=PP_ALIGN.CENTER)
            add_textbox(slide, Inches(8.34), Inches(5.02), Inches(3.1), Inches(0.58), spec["note"], size=13.5, color=MUTED)
        elif layout == "grid4":
            if spec["title"] == "为什么小程序样式开发一直慢":
                add_textbox(slide, Inches(9.95), Inches(1.78), Inches(1.1), Inches(0.7), "4", size=34, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
                add_textbox(slide, Inches(9.78), Inches(2.32), Inches(1.5), Inches(0.2), "个固定摩擦点", size=11.5, color=MUTED, align=PP_ALIGN.CENTER)
                add_textbox(slide, Inches(0.95), Inches(1.68), Inches(8.8), Inches(0.95), "真正拖慢你的，\n不是 CSS 本身。", size=28, bold=True)
                add_textbox(slide, Inches(0.98), Inches(2.62), Inches(5.8), Inches(0.2), "而是表达、试错、适配、反馈这 4 个摩擦点叠在一起。", size=11.8, color=MUTED)
                positions = [(0.95, 3.12), (6.42, 3.12), (0.95, 4.74), (6.42, 4.74)]
                colors = [PANEL_GRADIENT, PANEL_GRADIENT_ALT, PANEL_GRADIENT, ACCENT_GRADIENT]
                for (x, y), item, fill in zip(positions, spec["items"], colors):
                    add_card(slide, Inches(x), Inches(y), Inches(4.72), Inches(1.16), fill)
                    add_textbox(slide, Inches(x + 0.22), Inches(y + 0.34), Inches(4.15), Inches(0.28), item, size=18, bold=True)
                add_card(slide, Inches(11.55), Inches(2.1), Inches(0.34), Inches(3.95), ACCENT_GRADIENT, radius=MSO_AUTO_SHAPE_TYPE.RECTANGLE)
            elif spec["title"] == "为什么是 Tailwind，不是直接写 CSS":
                add_card(slide, Inches(0.95), Inches(2.0), Inches(4.25), Inches(3.95), PANEL_GRADIENT)
                add_textbox(slide, Inches(1.28), Inches(2.45), Inches(3.35), Inches(0.5), "结构化语言", size=28, color=CYAN, bold=True)
                add_textbox(slide, Inches(1.28), Inches(3.18), Inches(3.35), Inches(1.35), "比自由 CSS 更可枚举、更可约束，也更适合 AI 稳定输出。", size=16)
                positions = [(5.62, 2.05), (8.92, 2.05), (5.62, 4.08), (8.92, 4.08)]
                colors = [PANEL_GRADIENT, PANEL_GRADIENT_ALT, PANEL_GRADIENT, ACCENT_GRADIENT]
                for (x, y), item, fill in zip(positions, spec["items"], colors):
                    add_card(slide, Inches(x), Inches(y), Inches(2.95), Inches(1.58), fill)
                    add_textbox(slide, Inches(x), Inches(y + 0.55), Inches(2.95), Inches(0.25), item, size=18, bold=True, align=PP_ALIGN.CENTER)
            else:
                positions = [(0.9, 1.9), (6.75, 1.9), (0.9, 4.1), (6.75, 4.1)]
                colors = [PANEL_GRADIENT, PANEL_GRADIENT_ALT, PANEL_GRADIENT, ACCENT_GRADIENT]
                for (x, y), item, fill in zip(positions, spec["items"], colors):
                    add_card(slide, Inches(x), Inches(y), Inches(5.2), Inches(1.55), fill)
                    add_textbox(slide, Inches(x + 0.28), Inches(y + 0.46), Inches(4.6), Inches(0.5), item, size=21, bold=True)
        elif layout == "compare":
            add_card(slide, Inches(0.92), Inches(1.85), Inches(5.75), Inches(4.45), PANEL_GRADIENT)
            add_card(slide, Inches(6.78), Inches(1.85), Inches(5.65), Inches(4.45), PANEL_GRADIENT_ALT)
            add_textbox(slide, Inches(1.18), Inches(2.15), Inches(1.8), Inches(0.3), spec["left_title"], size=16, color=ROSE, bold=True)
            add_textbox(slide, Inches(7.05), Inches(2.15), Inches(1.8), Inches(0.3), spec["right_title"], size=16, color=EMERALD, bold=True)
            for i, item in enumerate(spec["left_items"]):
                add_textbox(slide, Inches(1.18), Inches(2.8 + i * 0.92), Inches(4.8), Inches(0.4), f"• {item}", size=21)
            for i, item in enumerate(spec["right_items"]):
                add_textbox(slide, Inches(7.05), Inches(2.8 + i * 0.92), Inches(4.8), Inches(0.4), f"• {item}", size=21)
            add_card(slide, Inches(5.98), Inches(2.2), Inches(0.1), Inches(3.55), ACCENT_GRADIENT, radius=MSO_AUTO_SHAPE_TYPE.RECTANGLE)
            add_card(slide, Inches(3.65), Inches(5.72), Inches(5.05), Inches(0.42), ACCENT_GRADIENT)
            add_textbox(slide, Inches(3.8), Inches(5.83), Inches(4.7), Inches(0.16), "核心转折：以前写的是属性，现在写的是意图。", size=12.5, bold=True, align=PP_ALIGN.CENTER)
            add_textbox(slide, Inches(4.24), Inches(1.34), Inches(4.8), Inches(0.24), "这页要停顿一下再讲。", size=11.2, color=MUTED, align=PP_ALIGN.CENTER)
        elif layout == "prompt":
            add_code_block(slide, Inches(0.9), Inches(1.9), Inches(6.15), Inches(4.65), "Prompt 示例", PROMPT_SNIPPET)
            card = add_card(slide, Inches(7.48), Inches(2.1), Inches(4.82), Inches(3.88), ACCENT_GRADIENT)
            card.line.color.rgb = CYAN
            add_textbox(slide, Inches(7.83), Inches(2.52), Inches(4.0), Inches(0.4), "渐变卡片", size=21, bold=True)
            add_textbox(slide, Inches(7.83), Inches(3.05), Inches(3.8), Inches(0.35), "AI 直接生成第一版 UI 语言", size=13, color=WHITE)
            for i, bullet in enumerate(spec["bullets"]):
                add_card(slide, Inches(7.82), Inches(3.74 + i * 0.48), Inches(1.66), Inches(0.33), PANEL_GRADIENT_ALT if i % 2 == 0 else ACCENT_GRADIENT)
                add_textbox(slide, Inches(8.02), Inches(3.8 + i * 0.48), Inches(3.4), Inches(0.2), bullet, size=11, color=TEXT, bold=True)
            add_card(slide, Inches(0.95), Inches(6.03), Inches(5.7), Inches(0.34), PANEL_GRADIENT_ALT)
            add_textbox(slide, Inches(1.12), Inches(6.12), Inches(5.35), Inches(0.14), "讲法提示：别先讲 Prompt 技巧，先让观众看到“结果真的出来了”。", size=11.5, color=MUTED)
            for i, label in enumerate(["描述界面", "拿到类名", "直接贴进页面"]):
                left = Inches(7.72 + i * 1.48)
                add_card(slide, left, Inches(5.78), Inches(1.28), Inches(0.3), PANEL_GRADIENT_ALT)
                add_textbox(slide, left, Inches(5.86), Inches(1.28), Inches(0.12), label, size=9.5, color=TEXT, align=PP_ALIGN.CENTER)
            add_textbox(slide, Inches(7.76), Inches(1.54), Inches(4.0), Inches(0.22), "先看结果，再回头拆过程。", size=11.5, color=CYAN, bold=True)
        elif layout == "flow3":
            colors = [PANEL_GRADIENT_ALT, ACCENT_GRADIENT, PANEL_GRADIENT_ALT]
            add_card(slide, Inches(0.95), Inches(1.95), Inches(3.1), Inches(0.56), PANEL_GRADIENT)
            add_textbox(slide, Inches(1.14), Inches(2.1), Inches(2.7), Inches(0.18), "记住这一张，整场直播都在解释它", size=12.5, color=MUTED)
            for i, item in enumerate(spec["items"]):
                left = Inches(0.95 + i * 4.0)
                add_card(slide, left, Inches(2.65), Inches(3.0), Inches(1.8), colors[i])
                add_textbox(slide, left + Inches(0.22), Inches(3.02), Inches(2.55), Inches(0.82), item, size=17, color=TEXT, bold=True)
                if i < 2:
                    add_textbox(slide, left + Inches(3.1), Inches(3.15), Inches(0.7), Inches(0.3), "→", size=24, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
            add_textbox(slide, Inches(0.95), Inches(5.65), Inches(11.0), Inches(0.35), spec["note"], size=14, color=MUTED)
            add_card(slide, Inches(8.82), Inches(1.95), Inches(2.88), Inches(0.48), PANEL_GRADIENT_ALT)
            add_textbox(slide, Inches(8.98), Inches(2.08), Inches(2.55), Inches(0.16), "先对接 Tailwind，再对接小程序", size=11.5, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
            add_textbox(slide, Inches(0.98), Inches(6.1), Inches(5.2), Inches(0.14), "你讲完这一页，观众应该已经知道整条链路怎么记。", size=11.2, color=MUTED)
        elif layout == "stack":
            add_panel_frame(slide, Inches(7.25), Inches(1.62), Inches(4.95), Inches(3.76), "Architecture")
            add_picture(slide, LAYER_IMG, Inches(7.52), Inches(1.98), width=Inches(4.4), height=Inches(3.1))
            add_tag(slide, Inches(0.95), Inches(1.82), Inches(2.05), "不是插件，是链路")
            widths = [5.6, 5.1, 4.6, 4.1]
            fills = [PANEL_GRADIENT, PANEL_GRADIENT_ALT, PANEL_GRADIENT, ACCENT_GRADIENT]
            for i, item in enumerate(spec["items"]):
                left = Inches(0.95 + i * 0.26)
                top = Inches(5.38 - i * 0.73)
                add_card(slide, left, top, Inches(widths[i]), Inches(0.6), fills[i])
                add_textbox(slide, left + Inches(0.18), top + Inches(0.14), Inches(widths[i] - 0.32), Inches(0.2), item, size=12.5, bold=True)
        elif layout == "chips":
            if spec["title"] == "支持范围":
                add_picture(slide, PLUGINS_IMG, Inches(0.95), Inches(2.02), width=Inches(5.45), height=Inches(1.95))
                add_picture(slide, FRAMEWORKS_IMG, Inches(6.82), Inches(2.02), width=Inches(5.45), height=Inches(1.95))
                for i, item in enumerate(spec["items"]):
                    left = Inches(0.98 + i * 1.95)
                    top = Inches(4.72)
                    add_card(slide, left, top, Inches(1.62), Inches(0.52), PANEL_GRADIENT_ALT)
                    add_textbox(slide, left, top + Inches(0.12), Inches(1.62), Inches(0.2), item, size=11, bold=True, align=PP_ALIGN.CENTER)
                add_textbox(slide, Inches(0.95), Inches(5.58), Inches(10.8), Inches(0.28), spec["note"], size=13.5, color=MUTED)
            elif spec["title"] == "Skill 先问什么":
                add_textbox(slide, Inches(0.95), Inches(1.82), Inches(6.4), Inches(0.52), "先问对，再输出。\n上下文不完整，AI 只会把错误放大。", size=22, bold=True)
                add_picture(slide, CREATE_PROJECT_IMG, Inches(8.1), Inches(1.88), width=Inches(3.55), height=Inches(1.72))
                for i, item in enumerate(spec["items"]):
                    left = Inches(0.98 + (i % 3) * 2.18)
                    top = Inches(3.45 + (i // 3) * 1.12)
                    add_card(slide, left, top, Inches(1.8), Inches(0.72), PANEL_GRADIENT if i % 2 == 0 else PANEL_GRADIENT_ALT)
                    add_textbox(slide, left, top + Inches(0.19), Inches(1.8), Inches(0.22), item, size=14.5, bold=True, align=PP_ALIGN.CENTER)
                add_textbox(slide, Inches(8.08), Inches(4.02), Inches(3.35), Inches(0.8), spec["note"], size=13.5, color=MUTED)
            else:
                for i, item in enumerate(spec["items"]):
                    left = Inches(0.95 + (i % 3) * 4.03)
                    top = Inches(2.3 + (i // 3) * 1.28)
                    add_card(slide, left, top, Inches(3.25), Inches(0.82), PANEL_GRADIENT_ALT)
                    add_textbox(slide, left + Inches(0.2), top + Inches(0.22), Inches(2.85), Inches(0.3), item, size=18, bold=True, align=PP_ALIGN.CENTER)
                add_textbox(slide, Inches(0.95), Inches(5.55), Inches(10.8), Inches(0.3), spec["note"], size=14, color=MUTED)
        elif layout == "repo":
            if spec["title"] == "仓库不是 PPT 工程":
                add_panel_frame(slide, Inches(7.02), Inches(0.92), Inches(5.25), Inches(0.84), "Brand")
                add_picture(slide, LOGO_WIDE, Inches(7.34), Inches(1.1), width=Inches(4.55), height=Inches(0.46))
                add_panel_frame(slide, Inches(7.02), Inches(1.82), Inches(5.2), Inches(3.25), "Release")
                add_picture(slide, RELEASE_IMG, Inches(7.24), Inches(2.14), width=Inches(4.78), height=Inches(2.7))
                add_panel_frame(slide, Inches(7.02), Inches(5.12), Inches(5.2), Inches(1.44), "Showcase")
                add_picture(slide, SHOWCASE_COLLAGE, Inches(7.18), Inches(5.34), width=Inches(4.88), height=Inches(1.02))
                add_textbox(slide, Inches(7.22), Inches(6.5), Inches(4.6), Inches(0.28), "不是概念图，是真仓库、真模板、真发布节奏，也有真实落地案例。", size=11.5, color=MUTED)
                for i, item in enumerate(spec["items"]):
                    left = Inches(0.95 + (i % 1) * 5.95)
                    top = Inches(1.96 + i * 0.82)
                    add_card(slide, left, top, Inches(5.55), Inches(0.58), PANEL_GRADIENT if i % 2 == 0 else PANEL_GRADIENT_ALT)
                    add_textbox(slide, left + Inches(0.18), top + Inches(0.15), Inches(5.05), Inches(0.18), item, size=13, bold=True)
            elif spec["title"] == "进阶技巧":
                add_textbox(slide, Inches(0.95), Inches(1.82), Inches(5.7), Inches(0.5), "最容易翻车的不是语法。\n而是类名生成、枚举边界和端侧约束。", size=20, bold=True)
                add_code_block(slide, Inches(7.15), Inches(1.9), Inches(5.0), Inches(2.3), "真实写法锚点", INDEX_SNIPPET)
                for i, item in enumerate(spec["items"]):
                    left = Inches(0.95 + (i % 1) * 5.95)
                    top = Inches(2.78 + i * 0.64)
                    add_card(slide, left, top, Inches(5.45), Inches(0.54), PANEL_GRADIENT if i % 2 == 0 else PANEL_GRADIENT_ALT)
                    add_textbox(slide, left + Inches(0.18), top + Inches(0.14), Inches(5.0), Inches(0.16), item, size=12.5, bold=True)
            else:
                for i, item in enumerate(spec["items"]):
                    left = Inches(0.95 + (i % 2) * 5.95)
                    top = Inches(1.95 + (i // 2) * 1.18)
                    add_card(slide, left, top, Inches(5.2), Inches(0.84), PANEL_GRADIENT if i % 2 == 0 else PANEL_GRADIENT_ALT)
                    add_textbox(slide, left + Inches(0.22), top + Inches(0.2), Inches(4.7), Inches(0.3), item, size=18, bold=True)
        elif layout == "focus":
            add_panel_frame(slide, Inches(0.95), Inches(1.86), Inches(5.9), Inches(4.2), "Demo Scope")
            add_textbox(slide, Inches(1.22), Inches(2.35), Inches(5.0), Inches(0.35), "主讲目录", size=13, color=CYAN, bold=True)
            add_textbox(slide, Inches(1.22), Inches(2.78), Inches(4.9), Inches(1.0), "demo/uni-app-tailwindcss-v4", size=24, bold=True)
            add_textbox(slide, Inches(1.22), Inches(4.6), Inches(4.7), Inches(0.5), spec["note"], size=15)
            add_code_block(slide, Inches(7.02), Inches(1.86), Inches(5.2), Inches(2.4), "index.vue 真实片段", INDEX_SNIPPET)
            add_picture(slide, DEMO_LOGO, Inches(7.18), Inches(4.48), width=Inches(0.82), height=Inches(0.82))
            add_footer_note(slide, Inches(1.2), Inches(5.32), Inches(4.9), "演示顺序：patch -> vite 插件 -> main.css -> 页面类名")
            for i, item in enumerate(spec["items"]):
                left = Inches(7.18 + (i % 2) * 2.4)
                top = Inches(5.0 + (i // 2) * 0.86)
                add_card(slide, left, top, Inches(2.18), Inches(0.95), ACCENT_GRADIENT if i < 2 else PANEL_GRADIENT_ALT)
                add_textbox(slide, left, top + Inches(0.3), Inches(2.18), Inches(0.3), item, size=18, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
        elif layout == "config":
            add_panel_frame(slide, Inches(0.95), Inches(1.72), Inches(11.45), Inches(4.72), "Config Panel")
            add_card(slide, Inches(1.08), Inches(1.98), Inches(11.18), Inches(0.72), ACCENT_GRADIENT)
            label = "真实配置片段" if spec["title"].endswith("1") else "插件链路片段" if spec["title"].endswith("2") else "样式入口片段"
            add_textbox(slide, Inches(1.36), Inches(2.18), Inches(10.6), Inches(0.25), label, size=19, color=TEXT, bold=True)
            snippet = POSTINSTALL_SNIPPET if spec["title"].endswith("1") else VITE_SNIPPET if spec["title"].endswith("2") else MAIN_CSS_SNIPPET
            code_title = "package.json" if spec["title"].endswith("1") else "vite.config.ts" if spec["title"].endswith("2") else "src/main.css"
            add_code_block(slide, Inches(1.08), Inches(2.96), Inches(6.52), Inches(3.0), code_title, snippet)
            if spec["title"].endswith("3"):
                add_picture(slide, LOGO, Inches(10.55), Inches(2.2), width=Inches(1.0), height=Inches(1.0))
            for i, item in enumerate(spec["items"]):
                left = Inches(7.9)
                top = Inches(3.02 + i * 0.94)
                add_card(slide, left, top, Inches(3.96), Inches(0.62), PANEL_GRADIENT_ALT if i != 1 else ACCENT_GRADIENT)
                add_textbox(slide, left + Inches(0.18), top + Inches(0.16), Inches(3.58), Inches(0.2), item, size=13.5, bold=True)
            reason = "很多坑其实不是 Tailwind 本身，而是 patch 没接好。" if spec["title"].endswith("1") else "Tailwind 生成后，还要再过一层小程序适配。" if spec["title"].endswith("2") else "这不是业务 CSS 杂物间，而是原子类入口。"
            add_footer_note(slide, Inches(7.92), Inches(6.02), Inches(3.95), reason)
        elif layout == "flow4":
            colors = [PANEL_GRADIENT_ALT, ACCENT_GRADIENT, PANEL_GRADIENT_ALT, ACCENT_GRADIENT]
            add_panel_frame(slide, Inches(0.95), Inches(1.72), Inches(4.42), Inches(1.9), "Skill Input")
            add_code_block(slide, Inches(1.08), Inches(1.98), Inches(4.12), Inches(1.42), "Skill 输入风格", SKILL_PROMPT_SNIPPET)
            add_panel_frame(slide, Inches(7.94), Inches(1.72), Inches(4.28), Inches(0.86), "Brand")
            add_picture(slide, LOGO_WIDE, Inches(8.18), Inches(1.95), width=Inches(3.7), height=Inches(0.44))
            for i, item in enumerate(spec["items"]):
                left = Inches(0.75 + i * 3.1)
                add_card(slide, left, Inches(4.05), Inches(2.35), Inches(1.02), colors[i])
                add_textbox(slide, left + Inches(0.14), Inches(4.31), Inches(2.05), Inches(0.4), item, size=15.5, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
                if i < 3:
                    add_textbox(slide, left + Inches(2.42), Inches(4.3), Inches(0.42), Inches(0.2), "→", size=24, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
            add_footer_note(slide, Inches(0.95), Inches(6.0), Inches(7.45), "好的 Skill 不只是更会写代码，而是更会按流程交付。")
        elif layout == "deliverables":
            add_panel_frame(slide, Inches(8.12), Inches(1.72), Inches(4.08), Inches(1.98), "Result Snapshot")
            add_picture(slide, CREATE_PROJECT_IMG, Inches(8.38), Inches(2.02), width=Inches(3.58), height=Inches(1.36))
            add_textbox(slide, Inches(0.95), Inches(1.92), Inches(6.4), Inches(0.52), "输出不是“建议”。输出应该是可复制、可验证、可回滚。", size=21, bold=True)
            for i, item in enumerate(spec["items"]):
                left = Inches(0.95 + (i % 3) * 4.0)
                top = Inches(2.82 + (i // 3) * 1.5)
                fill = [PANEL_GRADIENT, PANEL_GRADIENT_ALT, ACCENT_GRADIENT, PANEL_GRADIENT_ALT, ACCENT_GRADIENT][i]
                add_card(slide, left, top, Inches(3.2), Inches(1.05), fill)
                add_textbox(slide, left + Inches(0.16), top + Inches(0.32), Inches(2.88), Inches(0.32), item, size=17, bold=True, align=PP_ALIGN.CENTER)
            add_footer_note(slide, Inches(8.25), Inches(3.95), Inches(3.4), "团队价值不在回答，而在可落地结果。")
        elif layout == "metrics":
            hero = add_card(slide, Inches(0.95), Inches(2.1), Inches(3.6), Inches(3.2), PANEL_GRADIENT)
            add_textbox(slide, Inches(1.24), Inches(2.6), Inches(2.8), Inches(0.8), "可验证", size=31, color=CYAN, bold=True)
            add_textbox(slide, Inches(1.24), Inches(3.45), Inches(2.8), Inches(1.2), spec["note"], size=15)
            add_panel_frame(slide, Inches(4.82), Inches(1.86), Inches(3.52), Inches(2.34), "Release Signal")
            add_picture(slide, RELEASE_IMG, Inches(5.02), Inches(2.16), width=Inches(3.12), height=Inches(1.88))
            add_panel_frame(slide, Inches(8.32), Inches(1.86), Inches(3.86), Inches(2.34), "Runtime Layer")
            add_picture(slide, LAYER_IMG, Inches(8.52), Inches(2.16), width=Inches(3.46), height=Inches(1.88))
            add_panel_frame(slide, Inches(4.82), Inches(4.12), Inches(4.78), Inches(1.84), "Showcase Signal")
            add_picture(slide, SHOWCASE_COLLAGE, Inches(4.98), Inches(4.4), width=Inches(4.46), height=Inches(1.34))
            add_panel_frame(slide, Inches(1.15), Inches(4.92), Inches(3.1), Inches(1.02), "Signal")
            add_textbox(slide, Inches(1.42), Inches(5.22), Inches(2.55), Inches(0.36), "能跑，只是起点。\n可验证，才是工程信号。", size=13.5, bold=True)
            for i, item in enumerate(spec["items"]):
                left = Inches(9.85)
                top = Inches(4.22 + i * 0.72)
                add_card(slide, left, top, Inches(2.15), Inches(0.52), ACCENT_GRADIENT if i == 0 else PANEL_GRADIENT_ALT)
                add_textbox(slide, left, top + Inches(0.14), Inches(2.15), Inches(0.16), item, size=11.5, bold=True, align=PP_ALIGN.CENTER)
        elif layout == "closing":
            add_textbox(slide, Inches(0.9), Inches(1.55), Inches(7.2), Inches(1.0), spec["title"], size=30, bold=True)
            add_textbox(slide, Inches(0.92), Inches(2.5), Inches(5.2), Inches(0.4), spec["subtitle"], size=18, color=MUTED)
            add_panel_frame(slide, Inches(8.1), Inches(0.78), Inches(4.1), Inches(0.92), "Brand")
            add_picture(slide, LOGO_WIDE, Inches(8.35), Inches(1.02), width=Inches(3.6), height=Inches(0.58))
            for i, item in enumerate(spec["items"]):
                top = Inches(3.25 + i * 0.72)
                add_card(slide, Inches(0.95), top, Inches(6.35), Inches(0.48), PANEL_GRADIENT if i % 2 == 0 else PANEL_GRADIENT_ALT)
                add_textbox(slide, Inches(1.16), top + Inches(0.1), Inches(5.9), Inches(0.2), item, size=16, bold=True)
            add_panel_frame(slide, Inches(8.08), Inches(1.86), Inches(4.02), Inches(4.1), "Takeaway")
            add_card(slide, Inches(8.24), Inches(2.18), Inches(3.7), Inches(3.45), ACCENT_GRADIENT)
            add_textbox(slide, Inches(8.62), Inches(2.42), Inches(2.8), Inches(0.4), "Takeaway", size=14, color=TEXT, bold=True)
            add_textbox(slide, Inches(8.62), Inches(3.0), Inches(2.65), Inches(1.8), "别再从属性开始。\n先从意图开始。", size=24, color=TEXT, bold=True)


def main():
    ensure_generated_assets()
    if OUT.exists() and not BACKUP.exists():
        copy2(OUT, BACKUP)
    prs = Presentation()
    build(prs)
    prs.save(OUT)


if __name__ == "__main__":
    main()
